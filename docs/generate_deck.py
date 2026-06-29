#!/usr/bin/env python3
"""Generate SmartStudio capability presentation (.pptx) for engineering audience."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK = RGBColor(0x1E, 0x1B, 0x4B)       # indigo-950
PRIMARY = RGBColor(0x45, 0x38, 0xCA)     # indigo-600
ACCENT = RGBColor(0x06, 0xB6, 0xD4)      # cyan-500
SUCCESS = RGBColor(0x10, 0xB9, 0x81)     # emerald-500
WARN = RGBColor(0xF5, 0x9E, 0x0B)        # amber-500
GRAY = RGBColor(0x64, 0x74, 0x8B)        # slate-500
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)  # slate-100
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0F, 0x17, 0x2A)       # slate-900

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=16, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_metric_card(slide, left, top, value, label, color=PRIMARY):
    w, h = Inches(2.8), Inches(1.6)
    rect = add_rect(slide, left, top, w, h, LIGHT_GRAY)
    rect.shadow.inherit = False
    # accent bar at top
    add_rect(slide, left, top, w, Inches(0.06), color)
    add_text(slide, left + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.7),
             value, size=32, color=color, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.95), w - Inches(0.6), Inches(0.5),
             label, size=13, color=GRAY)

def add_icon_bullet(slide, left, top, width, icon, title, desc, color=PRIMARY):
    # icon circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, left + Inches(0.55), top, width - Inches(0.55), Inches(0.3),
             title, size=16, color=BLACK, bold=True)
    add_text(slide, left + Inches(0.55), top + Inches(0.3), width - Inches(0.55), Inches(0.5),
             desc, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.2),
         "SmartStudio", size=54, color=WHITE, bold=True)
add_text(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.8),
         "Metadata-Driven Full-Stack App Generation Platform", size=26, color=ACCENT)
add_text(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.6),
         "Design  \u2192  Configure  \u2192  Generate  \u2192  Deploy", size=18, color=RGBColor(0x94, 0xA3, 0xB8))

add_rect(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(0.01), RGBColor(0x33, 0x30, 0x6B))

add_text(slide, Inches(1.2), Inches(5.5), Inches(6), Inches(0.4),
         "Engineering Deep-Dive  |  Architecture  \u00B7  Performance  \u00B7  Code Generation", size=14, color=RGBColor(0x94, 0xA3, 0xB8))
add_text(slide, Inches(1.2), Inches(6.1), Inches(6), Inches(0.4),
         "March 2026", size=13, color=RGBColor(0x64, 0x74, 0x8B))


# ═══════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "The Problem", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "Every retail client needs an inventory management app. Every app is 80% the same.", size=18, color=GRAY)

# Pain point cards
cards = [
    ("6\u201312 months", "Time to build one app\nfrom scratch", "\u23F1"),
    ("48 DataViews", "Per app \u2014 each needs schema,\npipeline, filters, viewport", "\u26C5"),
    ("Rust + React", "gRPC backend + React\nfrontend per client", "\u2699"),
    ("80% repeated", "Same patterns, same pipeline\ninfra, same UI components", "\u267B"),
]

for i, (val, desc, icon) in enumerate(cards):
    left = Inches(0.8 + i * 3.1)
    top = Inches(2.0)
    w, h = Inches(2.8), Inches(2.8)
    rect = add_rect(slide, left, top, w, h, LIGHT_GRAY)
    add_text(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.6), Inches(0.6),
             icon, size=28, color=PRIMARY)
    add_text(slide, left + Inches(0.3), top + Inches(0.9), w - Inches(0.6), Inches(0.6),
             val, size=24, color=PRIMARY, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.6), w - Inches(0.6), Inches(1.0),
             desc, size=13, color=GRAY)

add_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2), RGBColor(0xEE, 0xF2, 0xFF))
add_text(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.8),
         "What if we could capture the app definition as metadata and generate the entire stack?",
         size=20, color=PRIMARY, bold=True)


# ═══════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "Architecture Overview", size=32, color=DARK, bold=True)

# Top row (frontend)
layers = [
    (Inches(0.8), Inches(1.5), Inches(5.5), "SmartStudio Web UI", "React 19  \u00B7  TanStack  \u00B7  Zustand  \u00B7  Tailwind CSS", PRIMARY),
    (Inches(6.8), Inches(1.5), Inches(5.5), "Frontend Preview (WYSIWYG)", "Live preview of generated app layout", ACCENT),
]
for left, top, w, title, sub, color in layers:
    h = Inches(1.0)
    add_rect(slide, left, top, w, h, color)
    add_text(slide, left + Inches(0.3), top + Inches(0.12), w - Inches(0.6), Inches(0.35),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.5), w - Inches(0.6), Inches(0.35),
             sub, size=11, color=RGBColor(0xE0, 0xE7, 0xFF))

# Arrow
add_text(slide, Inches(6), Inches(2.7), Inches(1), Inches(0.4), "\u2193", size=20, color=GRAY, align=PP_ALIGN.CENTER)

# Middle row (backend)
add_rect(slide, Inches(0.8), Inches(3.0), Inches(3.5), Inches(1.2), RGBColor(0xF0, 0xFD, 0xF4))
add_text(slide, Inches(1.0), Inches(3.1), Inches(3), Inches(0.3),
         "Rust Axum Server", size=15, color=RGBColor(0x16, 0x65, 0x34), bold=True)
add_text(slide, Inches(1.0), Inches(3.45), Inches(3), Inches(0.6),
         "REST API  \u00B7  SSE Streaming\nPipeline Executor  \u00B7  DuckDB Engine", size=11, color=GRAY)

add_rect(slide, Inches(4.6), Inches(3.0), Inches(3.0), Inches(1.2), RGBColor(0xFF, 0xF7, 0xED))
add_text(slide, Inches(4.8), Inches(3.1), Inches(2.6), Inches(0.3),
         "SQLite Metadata Store", size=15, color=RGBColor(0x9A, 0x34, 0x12), bold=True)
add_text(slide, Inches(4.8), Inches(3.45), Inches(2.6), Inches(0.6),
         "Apps  \u00B7  DataViews  \u00B7  Dimensions\nModules  \u00B7  Connections  \u00B7  Configs", size=11, color=GRAY)

add_rect(slide, Inches(7.9), Inches(3.0), Inches(4.4), Inches(1.2), RGBColor(0xEF, 0xF6, 0xFF))
add_text(slide, Inches(8.1), Inches(3.1), Inches(4), Inches(0.3),
         "Data Pipeline Engine", size=15, color=RGBColor(0x1E, 0x40, 0xAF), bold=True)
add_text(slide, Inches(8.1), Inches(3.45), Inches(4), Inches(0.6),
         "PG COPY \u2192 Parquet \u2192 DuckDB \u2192 Hive-partitioned output\nCDC  \u00B7  Parallel extraction  \u00B7  3-phase orchestrator", size=11, color=GRAY)

# Bottom row (outputs)
outputs = [
    ("Rust gRPC Backend", "Proto + Axum services"),
    ("React Frontend", "Module/Component scaffold"),
    ("Extensible Targets", "Language pack per platform"),
    ("Parquet Data Lake", "Hive-partitioned files"),
]
for i, (title, sub) in enumerate(outputs):
    left = Inches(0.8 + i * 3.1)
    add_rect(slide, left, Inches(4.8), Inches(2.8), Inches(0.9), DARK)
    add_text(slide, left + Inches(0.2), Inches(4.9), Inches(2.4), Inches(0.3),
             title, size=13, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.2), Inches(5.2), Inches(2.4), Inches(0.3),
             sub, size=10, color=RGBColor(0x94, 0xA3, 0xB8))

# Arrow between rows
add_text(slide, Inches(6), Inches(4.4), Inches(1), Inches(0.4), "\u2193  Code Gen  \u2193", size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
         "Single metadata definition \u2192 Multiple generated targets  |  One source of truth for schema, pipeline, UI, and filters",
         size=13, color=GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 4 — Data Pipeline: 3-Phase Orchestrator
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "3-Phase App Pipeline Orchestrator", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Extracts 26 PG source tables, deduplicates, materializes 48 DataViews \u2014 all via SSE-streamed progress", size=16, color=GRAY)

# Phase boxes
phases = [
    ("P1", "Extract", "Parallel PG COPY \u2192 Parquet",
     ["37 unique source tables", "Deduplicated by output_path", "Parallel OS threads", "CDC incremental support", "tokio multi-thread runtime"],
     SUCCESS),
    ("P2", "Materialize", "DuckDB Load \u2192 Join \u2192 Write",
     ["Load parquet into DuckDB", "Run join/aggregation queries", "Write hive-partitioned output", "Per-DV executor (isolation)", "Dependency DAG resolution"],
     PRIMARY),
    ("P3", "BQ Export", "BigQuery \u2192 GCS \u2192 Local",
     ["EXPORT DATA to GCS", "gsutil parallel download", "Parquet format", "Project/dataset from config", "2-step materialization"],
     WARN),
]

for i, (tag, title, sub, bullets, color) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    top = Inches(1.8)
    w = Inches(3.7)
    h = Inches(4.8)

    add_rect(slide, left, top, w, h, LIGHT_GRAY)
    add_rect(slide, left, top, w, Inches(0.08), color)

    # Phase tag
    tag_shape = add_rect(slide, left + Inches(0.2), top + Inches(0.25), Inches(0.5), Inches(0.35), color)
    tf = tag_shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, left + Inches(0.85), top + Inches(0.22), Inches(2.5), Inches(0.35),
             title, size=20, color=BLACK, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.75), w - Inches(0.4), Inches(0.35),
             sub, size=12, color=GRAY)

    add_bullet_list(slide, left + Inches(0.3), top + Inches(1.3), w - Inches(0.6), Inches(3.2),
                    [f"\u2022  {b}" for b in bullets], size=12, color=BLACK, spacing=Pt(4))


# ═══════════════════════════════════════════════════════
# SLIDE 5 — Current Codebase vs SmartStudio Generated
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), SUCCESS)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "Hand-Built vs SmartStudio Generated", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Comparing the current inventory-smart-rust codebase with what SmartStudio generates from metadata", size=16, color=GRAY)

# ── Left: Current Hand-Built ──
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2), RGBColor(0xFE, 0xF2, 0xF2))
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.08), RGBColor(0xDC, 0x26, 0x26))
add_text(slide, Inches(1.1), Inches(1.85), Inches(5), Inches(0.35),
         "Current: Hand-Built  (inventory-smart-rust)", size=18, color=RGBColor(0xDC, 0x26, 0x26), bold=True)

# Scale numbers
add_text(slide, Inches(1.1), Inches(2.35), Inches(5), Inches(0.3),
         "Scale", size=13, color=RGBColor(0x7F, 0x1D, 0x1D), bold=True)
add_bullet_list(slide, Inches(1.3), Inches(2.65), Inches(5), Inches(1.2),
    ["87K lines Rust  +  98K lines TypeScript  =  185K LOC",
     "28 Rust crates  \u00B7  16 modules  \u00B7  115+ HTTP endpoints",
     "11 client configs hand-maintained per module",
     "30+ router files  \u00B7  16 service files  \u00B7  1 proto file"],
    size=11, color=RGBColor(0x7F, 0x1D, 0x1D), spacing=Pt(3))

# Productivity
add_text(slide, Inches(1.1), Inches(3.9), Inches(5), Inches(0.3),
         "Productivity", size=13, color=RGBColor(0x7F, 0x1D, 0x1D), bold=True)
add_bullet_list(slide, Inches(1.3), Inches(4.2), Inches(5), Inches(1.0),
    ["\u2718  Each DataView hand-coded: router, service, data, models, queries",
     "\u2718  Client config JSON duplicated across 16 modules",
     "\u2718  Schema changes cascade through multiple crates manually"],
    size=11, color=RGBColor(0x7F, 0x1D, 0x1D), spacing=Pt(3))

# Performance
add_text(slide, Inches(1.1), Inches(5.2), Inches(5), Inches(0.3),
         "Pipeline Performance", size=13, color=RGBColor(0x7F, 0x1D, 0x1D), bold=True)
add_bullet_list(slide, Inches(1.3), Inches(5.5), Inches(5), Inches(1.0),
    ["\u2718  Sequential table extraction, one at a time",
     "\u2718  No shared extraction across DataViews",
     "\u2718  Full snapshot every refresh \u2014 no CDC incremental"],
    size=11, color=RGBColor(0x7F, 0x1D, 0x1D), spacing=Pt(3))

# ── Right: SmartStudio Generated ──
add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.2), RGBColor(0xF0, 0xFD, 0xF4))
add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.08), RGBColor(0x16, 0x65, 0x34))
add_text(slide, Inches(7.3), Inches(1.85), Inches(5), Inches(0.35),
         "SmartStudio Generated", size=18, color=RGBColor(0x16, 0x65, 0x34), bold=True)

# What it generates
add_text(slide, Inches(7.3), Inches(2.35), Inches(5), Inches(0.3),
         "Generated from Metadata", size=13, color=RGBColor(0x14, 0x53, 0x2D), bold=True)
add_bullet_list(slide, Inches(7.5), Inches(2.65), Inches(5), Inches(1.2),
    ["48 DataViews \u2192 48 gRPC services + REST + React UI",
     "Proto, handler, filter logic, column defs \u2014 all automatic",
     "One metadata change regenerates entire stack",
     "Language pack extensible to any target platform"],
    size=11, color=RGBColor(0x14, 0x53, 0x2D), spacing=Pt(3))

# Productivity
add_text(slide, Inches(7.3), Inches(3.9), Inches(5), Inches(0.3),
         "Productivity", size=13, color=RGBColor(0x14, 0x53, 0x2D), bold=True)
add_bullet_list(slide, Inches(7.5), Inches(4.2), Inches(5), Inches(1.0),
    ["\u2714  Define DataView metadata once in WYSIWYG editor",
     "\u2714  Schema + pipeline + filters + viewport in one place",
     "\u2714  New client = clone app + update connections"],
    size=11, color=RGBColor(0x14, 0x53, 0x2D), spacing=Pt(3))

# Performance
add_text(slide, Inches(7.3), Inches(5.2), Inches(5), Inches(0.3),
         "Pipeline Performance", size=13, color=RGBColor(0x14, 0x53, 0x2D), bold=True)
add_bullet_list(slide, Inches(7.5), Inches(5.5), Inches(5), Inches(1.0),
    ["\u2714  37 tables extracted in parallel (OS threads)",
     "\u2714  Shared extracts deduplicated across all 48 DVs",
     "\u2714  CDC incremental + SSE real-time progress"],
    size=11, color=RGBColor(0x14, 0x53, 0x2D), spacing=Pt(3))


# ═══════════════════════════════════════════════════════
# SLIDE 6 — Code Generation (Meta Capability)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "Code Generation: The Meta Capability", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "One metadata definition generates production-grade code across multiple targets", size=16, color=GRAY)

# Flow: Metadata → Language Pack → Generated Code
flow_items = [
    ("DataView\nMetadata", "48 DataViews with\nschema, pipeline,\nfilters, viewports", PRIMARY),
    ("\u2192", "", WHITE),
    ("Language\nPack", "Eta templates per\ntarget: Rust, React,\nextensible to any stack", ACCENT),
    ("\u2192", "", WHITE),
    ("Generated\nCode", "Production-grade\nscaffold: proto, gRPC,\nREST, UI components", SUCCESS),
]

for i, (title, desc, color) in enumerate(flow_items):
    if title == "\u2192":
        add_text(slide, Inches(1.0 + i * 2.3), Inches(2.4), Inches(1), Inches(0.5),
                 "\u2192", size=36, color=GRAY, align=PP_ALIGN.CENTER)
        continue
    left = Inches(0.8 + i * 2.3)
    w, h = Inches(2.5), Inches(2.2)
    add_rect(slide, left, Inches(1.8), w, h, color)
    add_text(slide, left + Inches(0.2), Inches(1.95), w - Inches(0.4), Inches(0.7),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.2), Inches(2.75), w - Inches(0.4), Inches(1.0),
             desc, size=11, color=RGBColor(0xE2, 0xE8, 0xF0))

# What gets generated
add_text(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.4),
         "What gets generated per DataView:", size=18, color=DARK, bold=True)

gen_items = [
    ("\u2699", "Proto Definition", "gRPC service with GetList,\nGetFilterValues RPCs", PRIMARY),
    ("\u2699", "Rust gRPC Service", "Axum handler + DuckDB\nparquet reader + filters", PRIMARY),
    ("\u2699", "REST Endpoints", "/list, /filter-values, /health\nwith pagination + search", ACCENT),
    ("\u2699", "React Component", "TanStack Table with column\ndefs, sort, search, filters", SUCCESS),
]

for i, (icon, title, desc, color) in enumerate(gen_items):
    add_icon_bullet(slide, Inches(0.8 + i * 3.1), Inches(5.0), Inches(2.8), icon, title, desc, color)


# ═══════════════════════════════════════════════════════
# SLIDE 7 — Visual Design / WYSIWYG
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "Visual Design: WYSIWYG Preview", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Design the app visually \u2014 see exactly what the generated app will look like", size=16, color=GRAY)

# Left panel — Design capabilities
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), LIGHT_GRAY)
add_text(slide, Inches(1.0), Inches(1.95), Inches(5), Inches(0.35),
         "Design Surface", size=18, color=DARK, bold=True)

design_features = [
    "\u2022  Module \u2192 SubModule \u2192 Component hierarchy",
    "\u2022  Drag-and-drop column reordering",
    "\u2022  Column visibility, sortable, searchable toggles",
    "\u2022  Filter configuration per dimension",
    "\u2022  Pipeline tree designer with node editor",
    "\u2022  Schema panel: types, groups, editability",
    "\u2022  ViewPort saved lens (filters + sort + pagination)",
    "\u2022  Inline SQL query preview with syntax highlighting",
]
add_bullet_list(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(4.0),
                design_features, size=13, color=BLACK, spacing=Pt(6))

# Right panel — Preview capabilities
add_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), LIGHT_GRAY)
add_text(slide, Inches(7.0), Inches(1.95), Inches(5), Inches(0.35),
         "Live Preview", size=18, color=DARK, bold=True)

preview_features = [
    "\u2022  Sidebar navigation (modules)",
    "\u2022  Tab layout (submodules \u2192 components)",
    "\u2022  Real data tables powered by DuckDB read_query",
    "\u2022  Per-column search (multi-search)",
    "\u2022  Pagination with configurable page sizes",
    "\u2022  Sort dropdown from DataView columns",
    "\u2022  Dimension filter panels",
    "\u2022  Mirrors generated app layout exactly",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(4.8), Inches(4.0),
                preview_features, size=13, color=BLACK, spacing=Pt(6))


# ═══════════════════════════════════════════════════════
# SLIDE 8 — Scalability
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "Scalability", size=32, color=DARK, bold=True)

# Dimension cards
dims = [
    ("Multi-Client", [
        "Per-client app definitions",
        "Client \u2192 Connections \u2192 Apps model",
        "Tenant-scoped config (TOML layers)",
        "Clone app across clients",
    ], PRIMARY),
    ("Multi-Environment", [
        "default \u2192 tenant.env \u2192 local.toml",
        "Merged config with override layers",
        "Per-environment pipeline settings",
        "BQ project / GCS bucket per env",
    ], ACCENT),
    ("Data Scale", [
        "Parquet + Hive partitioning",
        "DuckDB columnar query engine",
        "CDC incremental extraction",
        "Parallel I/O (OS threads)",
    ], SUCCESS),
    ("Code Scale", [
        "Language pack per target",
        "Template-driven (Eta engine)",
        "48 DataViews \u2192 48 services",
        "Proto + gRPC + REST + UI",
    ], WARN),
]

for i, (title, items, color) in enumerate(dims):
    left = Inches(0.8 + i * 3.1)
    top = Inches(1.5)
    w = Inches(2.8)

    add_rect(slide, left, top, w, Inches(4.5), LIGHT_GRAY)
    add_rect(slide, left, top, w, Inches(0.06), color)
    add_text(slide, left + Inches(0.25), top + Inches(0.2), w - Inches(0.5), Inches(0.4),
             title, size=18, color=color, bold=True)
    add_bullet_list(slide, left + Inches(0.25), top + Inches(0.8), w - Inches(0.5), Inches(3.5),
                    [f"\u2022  {b}" for b in items], size=12, color=BLACK, spacing=Pt(5))

# Bottom — numbers
add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), DARK)
nums = "48 DataViews  \u00B7  26 Source Tables  \u00B7  8 Modules  \u00B7  27 SubModules  \u00B7  33 Components  \u00B7  4 Filter Configs  \u00B7  3 Dimensions"
add_text(slide, Inches(1.2), Inches(6.45), Inches(11), Inches(0.5),
         nums, size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 9 — DataView Deep Dive
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "DataView: The Core Abstraction", size=32, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
         "Each DataView is a self-contained unit that defines schema, data pipeline, query interface, and UI behavior", size=16, color=GRAY)

# 5 sub-tabs
tabs = [
    ("Schema", "Column definitions: name, type,\nvisible, sortable, searchable,\neditable, filterable, group", PRIMARY),
    ("Pipeline", "Tree-based extraction workflow:\npg_extract \u2192 duckdb_table \u2192\nduckdb_query \u2192 loop", RGBColor(0x7C, 0x3A, 0xED)),
    ("Data", "Parquet browser + ad-hoc SQL\nconsole with {PARQUET_HOME}\nresolution", ACCENT),
    ("Filters", "Dimension bindings +\nViewPort CRUD (saved\nfilter+sort+pagination)", SUCCESS),
    ("ViewPorts", "Stateful filtered window\nwith read_query, sort,\npagination, column search", WARN),
]

for i, (title, desc, color) in enumerate(tabs):
    left = Inches(0.5 + i * 2.5)
    top = Inches(1.8)
    w, h = Inches(2.3), Inches(2.5)
    add_rect(slide, left, top, w, h, color)
    add_text(slide, left + Inches(0.2), top + Inches(0.2), w - Inches(0.4), Inches(0.4),
             title, size=18, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.7), w - Inches(0.4), Inches(1.5),
             desc, size=12, color=RGBColor(0xE2, 0xE8, 0xF0))

# Backend workflow
add_text(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.4),
         "Backend Workflow Pipeline (per DataView):", size=16, color=DARK, bold=True)

flow_steps = [
    ("Source", "pg_query\npg_sp\nbq_export"),
    ("\u2192", ""),
    ("Extract", "PG COPY\n\u2192 CSV\n\u2192 Parquet"),
    ("\u2192", ""),
    ("Transform", "DuckDB load\nJoin tables\nAggregate"),
    ("\u2192", ""),
    ("Output", "Hive-partitioned\nparquet files\n+ CDC marker"),
    ("\u2192", ""),
    ("Serve", "read_query\nDuckDB scan\nViewPort cache"),
]

for i, (title, desc) in enumerate(flow_steps):
    if title == "\u2192":
        add_text(slide, Inches(0.6 + i * 1.35), Inches(5.5), Inches(0.8), Inches(0.4),
                 "\u2192", size=24, color=GRAY, align=PP_ALIGN.CENTER)
        continue
    left = Inches(0.5 + i * 1.35)
    add_rect(slide, left, Inches(5.2), Inches(1.3), Inches(1.6), LIGHT_GRAY)
    add_text(slide, left + Inches(0.1), Inches(5.25), Inches(1.1), Inches(0.3),
             title, size=12, color=PRIMARY, bold=True)
    add_text(slide, left + Inches(0.1), Inches(5.55), Inches(1.1), Inches(1.1),
             desc, size=10, color=GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack Detail
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "Tech Stack", size=32, color=DARK, bold=True)

rows = [
    ("Frontend", "React 19, TanStack Table/Router, Zustand, Tailwind CSS, Lucide icons", PRIMARY),
    ("Backend", "Rust (Axum), SQLite (rusqlite), DuckDB, tokio-postgres, crossbeam", SUCCESS),
    ("Data Engine", "DuckDB in-memory columnar, Parquet (Snappy), Hive partitioning", ACCENT),
    ("Pipeline", "tokio multi-thread runtime, crossbeam channels, SSE streaming", RGBColor(0x7C, 0x3A, 0xED)),
    ("Code Gen", "Eta template engine, Language packs (Rust backend, React frontend)", WARN),
    ("Config", "TOML 3-layer override (default \u2192 tenant.env \u2192 local), SQLite traces", GRAY),
    ("External", "PostgreSQL (source), BigQuery (analytics), GCS (parquet lake)", RGBColor(0xDB, 0x27, 0x77)),
]

for i, (layer, tech, color) in enumerate(rows):
    top = Inches(1.3 + i * 0.78)
    add_rect(slide, Inches(0.8), top, Inches(2.2), Inches(0.6), color)
    add_text(slide, Inches(0.95), top + Inches(0.12), Inches(2), Inches(0.4),
             layer, size=15, color=WHITE, bold=True)
    add_text(slide, Inches(3.3), top + Inches(0.12), Inches(9), Inches(0.4),
             tech, size=14, color=BLACK)


# ═══════════════════════════════════════════════════════
# SLIDE 11 — What's Next
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "What's Next", size=32, color=DARK, bold=True)

roadmap = [
    ("Now", [
        "Fix 27 pipeline queries (PG schema alignment)",
        "Wire 8 stored procedure DataViews",
        "Dynamic ViewPort creation with cascading filters",
    ], SUCCESS),
    ("Next", [
        "Query abstraction in rust-shared-utils (QueryEngine trait)",
        "Schema introspection UI: browse PG \u2192 auto-populate columns",
        "Additional language pack templates",
        "Git integration for code generation (branch + commit)",
    ], PRIMARY),
    ("Future", [
        "Template marketplace (clone app definitions)",
        "Multi-tenant deployment orchestration",
        "Real-time CDC with Debezium/WAL",
        "AI-assisted pipeline builder",
    ], ACCENT),
]

for i, (phase, items, color) in enumerate(roadmap):
    left = Inches(0.8 + i * 4.1)
    top = Inches(1.5)
    w = Inches(3.7)

    add_rect(slide, left, top, w, Inches(4.5), LIGHT_GRAY)
    add_rect(slide, left, top, w, Inches(0.06), color)

    tag = add_rect(slide, left + Inches(0.2), top + Inches(0.25), Inches(0.8), Inches(0.35), color)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.9), w - Inches(0.6), Inches(3.3),
                    [f"\u2022  {b}" for b in items], size=13, color=BLACK, spacing=Pt(6))


# ═══════════════════════════════════════════════════════
# SLIDE 12 — Closing
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.0),
         "SmartStudio", size=48, color=WHITE, bold=True)
add_text(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.7),
         "Design once. Generate everywhere.", size=28, color=ACCENT)

add_rect(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.01), RGBColor(0x33, 0x30, 0x6B))

highlights = "Metadata-driven  \u00B7  48 DataViews  \u00B7  3-phase pipeline  \u00B7  37 parallel extracts  \u00B7  Multi-target code gen  \u00B7  WYSIWYG preview"
add_text(slide, Inches(1.2), Inches(4.9), Inches(11), Inches(0.5),
         highlights, size=16, color=RGBColor(0x94, 0xA3, 0xB8))

add_text(slide, Inches(1.2), Inches(5.8), Inches(6), Inches(0.4),
         "Questions?", size=22, color=WHITE, bold=True)


# Save
out_path = "/Users/karthickpachiappan/bb/smartstudio/docs/SmartStudio_Engineering_Deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
